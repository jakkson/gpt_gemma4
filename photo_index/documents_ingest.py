"""Ingest text from local documents (PDF, Office, plain text, etc.) into the shared FTS index."""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import sqlite3
import subprocess
import sys
import time
from datetime import datetime, timezone
from pathlib import Path

from photo_index.ingest_lock import global_ingest_lock
from photo_index.keep_awake import start_keep_awake
from photo_index.store import commit_ingest, connect, init_schema, upsert_photo


def _log(msg: str) -> None:
    print(msg, flush=True)


_DEFAULT_DB = Path(__file__).resolve().parent.parent / "data" / "photo_index.sqlite"


def documents_checkpoint_path(db_path: Path) -> Path:
    """Same directory as the SQLite DB, e.g. data/documents_ingest.checkpoint.json."""
    return db_path.with_name("documents_ingest.checkpoint.json")


def documents_skipped_path(db_path: Path) -> Path:
    """JSON report of files we couldn't index, e.g. data/documents_ingest.skipped.json."""
    return db_path.with_name("documents_ingest.skipped.json")


_SKIP_LIST_MAX = 10_000  # cap per-category lists so very large libraries don't OOM


def _fmt_elapsed(sec: float) -> str:
    if sec < 60:
        return f"{sec:.0f}s"
    m = int(sec // 60)
    s = int(sec % 60)
    if m < 60:
        return f"{m}m {s}s"
    h, m = m // 60, m % 60
    return f"{h}h {m}m"


def _write_documents_checkpoint(
    path: Path,
    *,
    root: str,
    db_path: Path,
    walk_files: int,
    candidates_seen: int,
    indexed: int,
    skipped_dup: int,
    skipped_empty: int,
    errors: int,
    skipped_video: int,
    skipped_hidden: int,
    skipped_media: int,
    skipped_temp: int,
    skipped_noise: int,
    last_rel_path: str,
    started_at_unix: float,
    finished: bool,
) -> None:
    now = time.time()
    elapsed = now - started_at_unix
    rate = (candidates_seen / elapsed) if elapsed > 0 else 0.0
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = {
        "root": root,
        "db_path": str(db_path.resolve()),
        "walk_files_seen": walk_files,
        "candidates_seen": candidates_seen,
        "indexed_this_run": indexed,
        "skipped_unchanged": skipped_dup,
        "skipped_empty": skipped_empty,
        "errors": errors,
        "skip_counts": {
            "video": skipped_video,
            "hidden": skipped_hidden,
            "image_or_audio": skipped_media,
            "temp_prefix": skipped_temp,
            "noise_dirs": skipped_noise,
        },
        "last_file_relpath": last_rel_path,
        "candidates_per_second": round(rate, 3),
        "elapsed_s": round(elapsed, 1),
        "elapsed_human": _fmt_elapsed(elapsed),
        "started_at_unix": started_at_unix,
        "updated_at_iso": datetime.now(timezone.utc).isoformat(),
        "ingest_finished": finished,
    }
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")

_VIDEO_EXT = frozenset(
    ".mp4 .m4v .mov .avi .mkv .webm .wmv .flv .mpg .mpeg .mts .m2ts .3gp .ogv".split()
)
_AUDIO_EXT = frozenset(".mp3 .m4a .aac .wav .aiff .aif .flac .ogg .wma .opus".split())
_RASTER_IMAGE_EXT = frozenset(
    ".jpg .jpeg .png .gif .bmp .tif .tiff .heic .webp .cr2 .nef .arw .dng .ico".split()
)
_SKIP_NAME_PREFIXES = ("~$", "._")

# Vendored / generated / cache directories — anything inside them is noise that
# pollutes FTS retrieval (e.g. setuptools' inflect lib brings in 130+ Norse-
# plural mentions of "valkyries", node_modules drowns out real content, etc.).
# Match exact path segments so we never index files whose path includes any of
# these names. Keep this in sync with NOISE_DIR_NAMES in prune_index.py.
NOISE_DIR_NAMES = frozenset({
    "site-packages",
    "node_modules",
    "__pycache__",
    ".pytest_cache",
    ".mypy_cache",
    ".ruff_cache",
    ".tox",
    ".turbo",
    ".next",
    ".cache",
    ".gradle",
    "bower_components",
    "vendor",
})

# Suffix-matched noise:
#   .egg-info, .dist-info: Python package metadata
#   .dcmd: Carbon Copy Cloner backup bundles (whole-Mac clones — not user
#          documents; one such bundle in this user's tree had >1.2M cached
#          icons inside it, which is the kind of pollution we never want
#          fed into FTS or VLM).
NOISE_DIR_SUFFIXES = (".egg-info", ".dist-info", ".dcmd")


def is_noise_path(rel_parts: tuple[str, ...]) -> bool:
    """True if any path segment matches a vendored/generated dir we never want indexed."""
    for seg in rel_parts:
        if seg in NOISE_DIR_NAMES:
            return True
        if seg.endswith(NOISE_DIR_SUFFIXES):
            return True
    return False


def doc_uuid(realpath_str: str) -> str:
    h = hashlib.sha256(realpath_str.encode("utf-8")).hexdigest()[:24]
    return f"doc:{h}"


def _meta_line(mtime_ns: int, size: int, rel: str, root_name: str) -> str:
    return (
        f"[docmeta] mtime_ns={mtime_ns} size={size}\n"
        f"source=dropbox_documents root={root_name}\nrel_path={rel}"
    )


def _should_skip_unchanged(
    conn: sqlite3.Connection,
    uuid: str,
    mtime_ns: int,
    size: int,
    force: bool,
) -> bool:
    if force:
        return False
    # A file is either one row (uuid) or many chunk rows (uuid#0, uuid#1, …);
    # check either representation for the unchanged mtime/size stamp.
    row = conn.execute(
        "SELECT vlm_text FROM photo_meta WHERE uuid IN (?, ?) LIMIT 1",
        (uuid, f"{uuid}#0"),
    ).fetchone()
    if row is None:
        return False
    blob = row[0] or ""
    m = re.search(r"mtime_ns=(\d+)\s+size=(\d+)", blob)
    if not m:
        return False
    return int(m.group(1)) == mtime_ns and int(m.group(2)) == size


# --- Chunking large documents (books, long PDFs) -----------------------------
#
# A whole book stored as ONE row is useless for RAG: the prompt only shows a
# small slice per record, so the model sees page 1 and nothing else. Split long
# text into overlapping passages, each its own indexed+embedded row, so the
# retriever can surface the RELEVANT passage and the model reads it in full.
_CHUNK_MIN_CHARS = 6_000     # below this: keep as a single row
_CHUNK_SIZE = 1_800          # target chars per chunk (fits the prompt field cap)
_CHUNK_OVERLAP = 250         # carry-over tail for cross-boundary context
_CHUNK_MAX = 600             # safety cap on chunks per file (very large books)

# Only PROSE files get chunked into passages. Data/code files (csv, json, xlsx,
# source code, plists…) stay a single truncated row: nobody retrieves "page 37
# of a contacts CSV", and chunking them would add ~100K noise rows competing
# with real content in every search. FTS still matches terms in the single row.
_PROSE_CHUNK_EXTS = frozenset({
    ".pdf", ".doc", ".docx", ".txt", ".rtf", ".md", ".markdown",
    ".html", ".htm", ".epub", ".pages", ".odt", ".pptx", ".key",
})


def _chunk_text(text: str) -> list[str]:
    """Split into ~_CHUNK_SIZE passages on paragraph boundaries, with overlap."""
    text = text.strip()
    if len(text) <= _CHUNK_MIN_CHARS:
        return [text]
    paras = re.split(r"\n\s*\n", text)
    chunks: list[str] = []
    cur = ""
    for p in paras:
        p = p.strip()
        if not p:
            continue
        if cur and len(cur) + len(p) + 2 > _CHUNK_SIZE:
            chunks.append(cur)
            cur = cur[-_CHUNK_OVERLAP:] + "\n\n" + p  # overlap tail for context
        else:
            cur = (cur + "\n\n" + p) if cur else p
        # A single paragraph larger than the target gets hard-split.
        while len(cur) > _CHUNK_SIZE * 1.6:
            chunks.append(cur[:_CHUNK_SIZE])
            cur = cur[_CHUNK_SIZE - _CHUNK_OVERLAP:]
        if len(chunks) >= _CHUNK_MAX:
            break
    if cur.strip() and len(chunks) < _CHUNK_MAX:
        chunks.append(cur)
    return chunks[:_CHUNK_MAX]


_TEXT_EXT = frozenset(
    """
    .txt .md .markdown .rst .csv .tsv .json .jsonl .xml .yaml .yml .ini .cfg .conf .log
    .html .htm .css .scss .sass .js .jsx .ts .tsx .mjs .cjs .py .rb .php .java .swift .kt
    .go .rs .c .cc .cpp .h .hpp .m .sql .sh .bash .zsh .plist .gradle
    """.split()
)


def _truncate(s: str, max_chars: int) -> str:
    if len(s) <= max_chars:
        return s
    return s[: max_chars - 80] + "\n\n… [truncated for index size]\n"


def extract_text_plain(path: Path) -> str:
    raw = path.read_bytes()
    if raw.startswith(b"\xff\xfe") or raw.startswith(b"\xfe\xff"):
        return raw.decode("utf-16", errors="replace")
    for enc in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return raw.decode(enc)
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace")


def extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path), strict=False)
    parts: list[str] = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(parts)


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    lines = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            lines.append("\t".join(c.text.strip() for c in row.cells))
    return "\n".join(lines)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                lines.append(shape.text_frame.text)
            elif hasattr(shape, "text"):
                tx = getattr(shape, "text", None)
                if tx:
                    lines.append(tx)
    return "\n".join(lines)


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in wb.worksheets:
            parts.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(max_row=5000, values_only=True):
                cells = []
                for c in row:
                    if c is None:
                        cells.append("")
                    else:
                        cells.append(str(c).strip())
                if any(x for x in cells):
                    parts.append("\t".join(cells))
    finally:
        wb.close()
    return "\n".join(parts)


def extract_xls(path: Path) -> str:
    import xlrd

    book = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for si in range(book.nsheets):
        sh = book.sheet_by_index(si)
        parts.append(f"## Sheet: {sh.name}")
        for rx in range(min(sh.nrows, 5000)):
            row = []
            for cx in range(sh.ncols):
                row.append(str(sh.cell_value(rx, cx)).strip())
            if any(row):
                parts.append("\t".join(row))
    return "\n".join(parts)


def extract_rtf(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text

    txt = extract_text_plain(path)
    return rtf_to_text(txt)


def extract_epub(path: Path) -> str:
    """EPUB = a zip of XHTML chapters. Stdlib-only: read the OPF spine for the
    author's chapter order (fallback: name-sorted), strip HTML to plain text
    with the same stripper the mail ingest uses."""
    import posixpath
    import urllib.parse
    import zipfile

    from photo_index.mail_ingest import _strip_html

    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        html_names = [n for n in names if n.lower().endswith((".xhtml", ".html", ".htm"))]
        ordered: list[str] | None = None
        opf = next((n for n in names if n.lower().endswith(".opf")), None)
        if opf:
            try:
                opf_xml = z.read(opf).decode("utf-8", "replace")
                manifest = dict(
                    re.findall(r'<item\b[^>]*\bid="([^"]+)"[^>]*\bhref="([^"]+)"', opf_xml)
                )
                # Some OPFs put href before id.
                for href, iid in re.findall(
                    r'<item\b[^>]*\bhref="([^"]+)"[^>]*\bid="([^"]+)"', opf_xml
                ):
                    manifest.setdefault(iid, href)
                spine = re.findall(r'<itemref\b[^>]*\bidref="([^"]+)"', opf_xml)
                base = posixpath.dirname(opf)
                cand = []
                for iid in spine:
                    href = manifest.get(iid)
                    if not href:
                        continue
                    p = posixpath.normpath(posixpath.join(base, urllib.parse.unquote(href)))
                    if p in names:
                        cand.append(p)
                if cand:
                    ordered = cand
            except Exception:
                ordered = None
        parts: list[str] = []
        for n in (ordered or sorted(html_names))[:500]:
            try:
                t = _strip_html(z.read(n).decode("utf-8", "replace"))
            except Exception:
                continue
            if t:
                parts.append(t)
    return "\n\n".join(parts)


def extract_textutil(path: Path) -> str | None:
    """macOS `textutil` converts legacy Word / Pages / RTF / HTML and more."""
    try:
        r = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True,
            text=True,
            timeout=300,
            check=False,
        )
    except Exception:
        return None
    if r.returncode != 0 or not (r.stdout or "").strip():
        return None
    return r.stdout


def extract_auto(path: Path, ext: str) -> tuple[str | None, str, str | None]:
    """Return (text, extractor_name, error_str_or_None)."""
    ext = ext.lower()
    err_str: str | None = None

    try:
        if ext == ".pdf":
            return extract_pdf(path), "pypdf", None
        if ext == ".docx":
            return extract_docx(path), "python-docx", None
        if ext == ".pptx":
            return extract_pptx(path), "python-pptx", None
        if ext == ".xlsx":
            return extract_xlsx(path), "openpyxl", None
        if ext == ".xls":
            return extract_xls(path), "xlrd", None
        if ext == ".epub":
            return extract_epub(path), "epub-zip", None
        if ext == ".rtf":
            try:
                return extract_rtf(path), "striprtf", None
            except Exception as e:
                err_str = f"{type(e).__name__}: {e}"
                t = extract_textutil(path)
                return (t if t else None), "striprtf|textutil", err_str
        if ext in _TEXT_EXT:
            return extract_text_plain(path), "plaintext", None
    except ImportError as e:
        err_str = f"ImportError: {e}"
        _log(f"[warn] Missing optional dependency for {path.name}: {e}")
    except Exception as e:
        # Placeholder/synced-but-not-materialized files (e.g. Dropbox), corrupt zips,
        # etc. Avoid dumping tracebacks into the log — fall through to textutil.
        err_str = f"{type(e).__name__}: {e}"
        _log(f"[documents extract warn] {path.name}: {err_str}")

    # Legacy Office / Pages / ambiguous types on macOS
    if sys.platform == "darwin":
        t = extract_textutil(path)
        if t and len(t.strip()) > 10:
            return t, "textutil", err_str

    return None, "unsupported", err_str


def run_documents_ingest(
    *,
    root: Path,
    index_db_path: Path,
    limit: int | None,
    force: bool,
    commit_every: int,
    progress_every: int,
    checkpoint_every: int,
    skip_images: bool,
    skip_audio: bool,
    max_chars_per_file: int,
    exclude: tuple[str, ...] = (),
) -> dict[str, int | float]:
    root = Path(os.path.abspath(root)).resolve()
    if not root.is_dir():
        raise NotADirectoryError(f"Root is not a directory: {root}")
    # Absolute, resolved prefixes to skip entirely (e.g. a noisy subfolder).
    exclude_prefixes = tuple(
        str(Path(os.path.abspath(os.path.expanduser(e))).resolve()) for e in exclude
    )
    if exclude_prefixes:
        _log(f"[documents] excluding paths under: {', '.join(exclude_prefixes)}")

    conn = connect(index_db_path)
    init_schema(conn)

    root_rp = root.resolve()
    root_display = root.name

    skipped_video = skipped_media = skipped_hidden = skipped_temp = 0
    skipped_noise = 0
    indexed = skipped_dup = skipped_empty = errors = 0
    candidate_i = 0
    walk_files = 0
    last_rel_path = ""
    t0 = time.perf_counter()
    started_at_unix = time.time()
    ck_path = documents_checkpoint_path(index_db_path)
    skipped_path = documents_skipped_path(index_db_path)

    empty_list: list[dict] = []
    error_list: list[dict] = []
    empty_overflow = 0
    error_overflow = 0

    lim_note = f" (limit {limit} files)" if limit is not None else ""
    _log(
        f"[documents] Streaming ingest under {root}{lim_note} — "
        "progress = candidate files (non-video, after filters). "
        f"Checkpoint JSON: {ck_path}"
    )
    _log(f"[documents] Skipped-files report will be written to: {skipped_path}")

    def rel_for(p_abs: Path) -> str:
        try:
            return str(p_abs.resolve().relative_to(root_rp))
        except Exception:
            return str(p_abs)

    def _write_skipped_report(finished: bool) -> None:
        try:
            payload = {
                "root": str(root),
                "db_path": str(index_db_path.resolve()),
                "started_at_unix": started_at_unix,
                "updated_at_iso": datetime.now(timezone.utc).isoformat(),
                "ingest_finished": finished,
                "summary": {
                    "walk_files_seen": walk_files,
                    "candidates_seen": candidate_i,
                    "indexed": indexed,
                    "skipped_unchanged": skipped_dup,
                    "skipped_empty_or_unreadable": skipped_empty,
                    "extract_errors": errors,
                    "filtered_videos": skipped_video,
                    "filtered_hidden": skipped_hidden,
                    "filtered_image_or_audio": skipped_media,
                    "filtered_temp_prefix": skipped_temp,
                    "filtered_noise_dirs": skipped_noise,
                },
                "list_caps": {
                    "max_per_list": _SKIP_LIST_MAX,
                    "empty_overflow_not_listed": empty_overflow,
                    "error_overflow_not_listed": error_overflow,
                },
                "skipped_empty_files": empty_list,
                "extract_error_files": error_list,
            }
            skipped_path.parent.mkdir(parents=True, exist_ok=True)
            skipped_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
        except Exception as e:
            _log(f"[documents skipped-report warn] could not write {skipped_path}: {e}")

    def _maybe_checkpoint(finished: bool = False) -> None:
        if finished:
            _write_documents_checkpoint(
                ck_path,
                root=str(root),
                db_path=index_db_path,
                walk_files=walk_files,
                candidates_seen=candidate_i,
                indexed=indexed,
                skipped_dup=skipped_dup,
                skipped_empty=skipped_empty,
                errors=errors,
                skipped_video=skipped_video,
                skipped_hidden=skipped_hidden,
                skipped_media=skipped_media,
                skipped_temp=skipped_temp,
                skipped_noise=skipped_noise,
                last_rel_path=last_rel_path,
                started_at_unix=started_at_unix,
                finished=True,
            )
            return
        if checkpoint_every <= 0:
            return
        if candidate_i == 0:
            return
        if candidate_i != 1 and (candidate_i % checkpoint_every != 0):
            return
        _write_documents_checkpoint(
            ck_path,
            root=str(root),
            db_path=index_db_path,
            walk_files=walk_files,
            candidates_seen=candidate_i,
            indexed=indexed,
            skipped_dup=skipped_dup,
            skipped_empty=skipped_empty,
            errors=errors,
            skipped_video=skipped_video,
            skipped_hidden=skipped_hidden,
            skipped_media=skipped_media,
            skipped_temp=skipped_temp,
            skipped_noise=skipped_noise,
            last_rel_path=last_rel_path,
            started_at_unix=started_at_unix,
            finished=False,
        )

    for p in root.rglob("*"):
        if not p.is_file():
            continue
        walk_files += 1
        if p.name.startswith(_SKIP_NAME_PREFIXES):
            skipped_temp += 1
            continue
        ext = p.suffix.lower()
        if ext in _VIDEO_EXT:
            skipped_video += 1
            continue
        if skip_audio and ext in _AUDIO_EXT:
            skipped_media += 1
            continue
        if skip_images and ext in _RASTER_IMAGE_EXT:
            skipped_media += 1
            continue
        try:
            rp = p.resolve()
            rel_parts = rp.relative_to(root_rp).parts
        except ValueError:
            continue
        if exclude_prefixes and str(rp).startswith(exclude_prefixes):
            skipped_noise += 1
            continue
        if any(seg.startswith(".") for seg in rel_parts):
            skipped_hidden += 1
            continue
        if is_noise_path(rel_parts):
            skipped_noise += 1
            continue

        candidate_i += 1
        path = p
        last_rel_path = rel_for(path)
        if limit is not None and candidate_i > limit:
            break

        elapsed = time.perf_counter() - t0
        rate = (candidate_i / elapsed) if elapsed > 0 else 0.0
        if progress_every and (candidate_i - 1) % progress_every == 0:
            tail = last_rel_path if len(last_rel_path) <= 100 else "…" + last_rel_path[-97:]
            _log(
                f"[documents progress] candidate {candidate_i} | walk_files={walk_files} | "
                f"indexed={indexed} dup={skipped_dup} empty={skipped_empty} err={errors} | "
                f"{rate:.2f} cand/s | elapsed {_fmt_elapsed(elapsed)} | {tail}"
            )

        _maybe_checkpoint(finished=False)
        if checkpoint_every > 0 and (candidate_i == 1 or candidate_i % checkpoint_every == 0):
            _write_skipped_report(finished=False)

        try:
            st = path.stat()
        except OSError:
            skipped_empty += 1
            continue

        real = os.path.realpath(str(path))
        uuid = doc_uuid(real)
        mtime_ns = int(getattr(st, "st_mtime_ns", int(st.st_mtime * 1_000_000_000)))
        size_i = int(st.st_size)

        if _should_skip_unchanged(conn, uuid, mtime_ns, size_i, force):
            skipped_dup += 1
            continue

        text_raw, how, extract_err = extract_auto(path, ext)
        rel_for_report = rel_for(path)
        if not text_raw or not text_raw.strip():
            skipped_empty += 1
            _log(f"[documents skip empty] {path.name} ({how})")
            if extract_err:
                # Couldn't extract anything AND an exception was raised → also
                # surface in the errors bucket so it shows up under "errors".
                errors += 1
                if len(error_list) < _SKIP_LIST_MAX:
                    error_list.append({
                        "rel_path": rel_for_report,
                        "ext": ext,
                        "extractor": how,
                        "error": extract_err,
                    })
                else:
                    error_overflow += 1
            else:
                if len(empty_list) < _SKIP_LIST_MAX:
                    empty_list.append({
                        "rel_path": rel_for_report,
                        "ext": ext,
                        "extractor": how,
                    })
                else:
                    empty_overflow += 1
            continue

        # Chunk long PROSE docs into passages; short files and data/code files
        # stay a single row. Cap total per file so a pathological doc can't
        # explode the index.
        if ext in _PROSE_CHUNK_EXTS:
            chunks = _chunk_text(text_raw.strip())
        else:
            chunks = [text_raw.strip()]
        if len(chunks) == 1:
            chunks = [_truncate(chunks[0], max_chars_per_file)]

        mt = getattr(st, "st_mtime", 0)
        iso = None
        if mt > 0:
            iso = datetime.fromtimestamp(mt, tz=timezone.utc).isoformat()

        rel = rel_for(path)
        meta = _meta_line(mtime_ns, size_i, rel, root_display)
        how_line = f"extractor={how}"

        # Clean slate for this file: remove any prior rows (a former single-row
        # book becoming chunks, or a changed chunk count) before writing.
        try:
            conn.execute(
                "DELETE FROM photo_meta WHERE uuid = ? OR uuid LIKE ?",
                (uuid, uuid + "#%"),
            )
        except Exception:
            pass

        n = len(chunks)
        wrote = 0
        for i, chunk in enumerate(chunks):
            cu = uuid if n == 1 else f"{uuid}#{i}"
            fn = rel if n == 1 else f"{rel} [part {i + 1}/{n}]"
            try:
                upsert_photo(
                    conn,
                    uuid=cu,
                    filename=fn,
                    date_iso=iso,
                    ocr_text=chunk,
                    vlm_text=f"{meta}\n{how_line}",
                    image_path_used=str(path),
                    commit=False,
                )
                wrote += 1
            except Exception as e:
                errors += 1
                _log(f"[documents warn] {path} chunk {i + 1}/{n}: {e}")
        if wrote:
            indexed += 1
            if commit_every <= 1 or indexed % commit_every == 0:
                commit_ingest(conn)

    commit_ingest(conn)
    _maybe_checkpoint(finished=True)
    _write_skipped_report(finished=True)
    elapsed = time.perf_counter() - t0
    _log(
        f"[documents done] candidates_seen={candidate_i} indexed={indexed} "
        f"skipped_unchanged={skipped_dup} skipped_empty={skipped_empty} errors={errors} "
        f"skip_counts: video≈{skipped_video} hidden≈{skipped_hidden} media≈{skipped_media} "
        f"temp≈{skipped_temp} noise_dirs≈{skipped_noise} "
        f"wall_time={elapsed:.1f}s db={index_db_path}"
    )
    _log(
        f"[documents] Skipped-files report: "
        f"empty={len(empty_list)} (+{empty_overflow} not listed) "
        f"errors={len(error_list)} (+{error_overflow} not listed) "
        f"→ {skipped_path}"
    )
    return {
        "indexed": indexed,
        "skipped_dup": skipped_dup,
        "skipped_empty": skipped_empty,
        "errors": errors,
        "elapsed": elapsed,
        "total": candidate_i,
    }


def main(argv: list[str] | None = None) -> None:
    p = argparse.ArgumentParser(
        description="Index PDF / Office / text files under a folder into photo_index.sqlite."
    )
    p.add_argument(
        "--root",
        type=str,
        default=str(Path.home() / "Dropbox" / "Documents"),
        help="Folder to walk recursively (default: ~/Dropbox/Documents).",
    )
    p.add_argument("--db", type=str, default=str(_DEFAULT_DB), help="Target SQLite DB path.")
    p.add_argument("--limit", type=int, default=None, help="Maximum number of files to try.")
    p.add_argument("--force", action="store_true", help="Re-read every file even if unchanged.")
    p.add_argument("--commit-every", type=int, default=25, help="Commit every N new rows.")
    p.add_argument(
        "--progress-every",
        type=int,
        default=200,
        help="Log a [documents progress] line every N candidate files (0=off).",
    )
    p.add_argument(
        "--checkpoint-every",
        type=int,
        default=25,
        help="Write documents_ingest.checkpoint.json every N candidates (0=only at end). Same folder as --db.",
    )
    p.add_argument(
        "--include-images",
        action="store_true",
        help="Also walk raster image files (normally skipped; photos belong in Photos ingest).",
    )
    p.add_argument(
        "--include-audio",
        action="store_true",
        help="Include audio extensions (normally skipped — little extractable text).",
    )
    p.add_argument(
        "--max-chars-per-file",
        type=int,
        default=400_000,
        help="Truncate stored text beyond this character count.",
    )
    p.add_argument(
        "--exclude",
        action="append",
        default=[],
        metavar="PATH",
        help="Skip files under this path (repeatable). E.g. a noisy subfolder.",
    )
    p.add_argument(
        "--no-global-ingest-lock",
        action="store_true",
        help="Don't use shared content-ingest.lock (not recommended).",
    )
    p.add_argument(
        "--no-keep-awake",
        action="store_true",
        help="Don't run caffeinate (long runs may sleep).",
    )
    args = p.parse_args(argv)

    index_db_path = Path(os.path.abspath(args.db))
    root = Path(os.path.abspath(args.root))

    if args.commit_every < 1:
        p.error("--commit-every must be >= 1")

    if not args.no_keep_awake:
        start_keep_awake(_log)

    def inner() -> None:
        run_documents_ingest(
            root=root,
            index_db_path=index_db_path,
            limit=args.limit,
            force=args.force,
            commit_every=args.commit_every,
            progress_every=args.progress_every,
            checkpoint_every=args.checkpoint_every,
            skip_images=not args.include_images,
            skip_audio=not args.include_audio,
            max_chars_per_file=max(10_000, args.max_chars_per_file),
            exclude=tuple(args.exclude),
        )

    if args.no_global_ingest_lock:
        inner()
        return

    with global_ingest_lock() as have_lock:
        if not have_lock:
            _log("[lock] Another content ingest is already running; skipping this run.")
            return
        inner()


if __name__ == "__main__":
    main()
